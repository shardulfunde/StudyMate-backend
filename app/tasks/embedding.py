from app.core.celery_app import celery
from app.db.session import SessionLocal
from app.db.models.resource import Resource
from app.db.models.resource_embedding import ResourceEmbedding
from app.utils.s3 import get_s3_client, get_bucket_name

import logging
import fitz
import io
from pptx import Presentation
from docx import Document
from dotenv import load_dotenv
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_cohere import CohereEmbeddings

load_dotenv()

S3_BUCKET = get_bucket_name()
logger = logging.getLogger(__name__)

BATCH_SIZE = 64


# -----------------------------
# Database Session
# -----------------------------
def create_db_session():
    return SessionLocal()


# -----------------------------
# File Extraction
# -----------------------------

def extract_text_from_pdf(file_bytes: bytes):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages = []

    for page_number, page in enumerate(doc, start=1):
        text = page.get_text()
        if text.strip():
            pages.append((page_number, text))

    doc.close()
    return pages


def extract_text_from_ppt(file_bytes: bytes):
    presentation = Presentation(io.BytesIO(file_bytes))
    slides = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)

        full_text = "\n".join(slide_text).strip()

        if full_text:
            slides.append((slide_number, full_text))

    return slides


def extract_text_from_docx(file_bytes: bytes):
    document = Document(io.BytesIO(file_bytes))
    full_text = []

    for para in document.paragraphs:
        if para.text.strip():
            full_text.append(para.text)

    combined_text = "\n".join(full_text).strip()

    if combined_text:
        return [(1, combined_text)]

    return []


def extract_text(file_bytes: bytes, file_key: str):
    file_key = file_key.lower()

    if file_key.endswith(".pdf"):
        return extract_text_from_pdf(file_bytes)

    elif file_key.endswith((".ppt", ".pptx")):
        return extract_text_from_ppt(file_bytes)

    elif file_key.endswith(".docx"):
        return extract_text_from_docx(file_bytes)

    else:
        raise ValueError("Unsupported file type")


# -----------------------------
# Batch Embedding Function
# -----------------------------
def batch_embed(model, texts, batch_size=BATCH_SIZE):
    all_vectors = []

    for i in range(0, len(texts), batch_size):
        batch = texts[i:i + batch_size]
        vectors = model.embed_documents(batch)
        all_vectors.extend(vectors)

    return all_vectors


# -----------------------------
# Celery Task
# -----------------------------
@celery.task
def generate_embeddings(resource_id: str):
    db = create_db_session()
    resource = None

    try:
        resource = db.query(Resource).filter(
            Resource.id == resource_id
        ).first()

        if not resource:
            logger.warning(
                "Embedding task skipped: resource %s not found",
                resource_id
            )
            return

        resource.embedding_status = "processing"
        db.commit()

        # Delete old embeddings (important for retries)
        db.query(ResourceEmbedding).filter(
            ResourceEmbedding.resource_id == resource.id
        ).delete()
        db.commit()

        # Download file
        s3 = get_s3_client()
        response = s3.get_object(
            Bucket=S3_BUCKET,
            Key=resource.file_key
        )

        file_bytes = response["Body"].read()

        pages = extract_text(file_bytes, resource.file_key)

        if not pages:
            logger.warning(
                "No extractable text found for resource %s",
                resource_id
            )
            resource.embedding_status = "failed"
            db.commit()
            return

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200
        )

        model = CohereEmbeddings(
            model="embed-english-v3.0"  # or embed-multilingual-v3.0
        )

        all_chunks = []
        chunk_metadata = []

        for page_number, text in pages:
            chunks = splitter.split_text(text)

            for idx, chunk in enumerate(chunks):
                all_chunks.append(chunk)
                chunk_metadata.append((page_number, idx))

        if not all_chunks:
            logger.warning(
                "No chunks generated for resource %s",
                resource_id
            )
            resource.embedding_status = "failed"
            db.commit()
            return

        # 🔹 Batched embedding
        vectors = batch_embed(model, all_chunks)

        for (chunk, vector), (page_number, idx) in zip(
            zip(all_chunks, vectors),
            chunk_metadata
        ):
            db.add(
                ResourceEmbedding(
                    resource_id=resource.id,
                    college_id=resource.college_id,
                    subject_id=resource.subject_id,
                    chunk_text=chunk,
                    page_number=page_number,
                    chunk_index=idx,
                    embedding=vector
                )
            )

        db.commit()

        resource.embedding_status = "completed"
        db.commit()

        logger.info(
            "Embeddings successfully generated for resource %s",
            resource_id
        )

    except Exception:
        logger.exception(
            "Embedding task failed for resource %s",
            resource_id
        )

        if resource is not None:
            resource.embedding_status = "failed"
            db.commit()

    finally:
        db.close()
